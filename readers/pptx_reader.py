"""
PowerPoint Reader for ThreadBear

Handles .pptx files using python-pptx with slide-based segmentation.
"""
from pathlib import Path
from .registry import reader_registry


class PptxReader:
    """Reader for PowerPoint presentations."""

    @staticmethod
    def extract_text(path):
        """Extract text from all slides."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("PowerPoint reading requires python-pptx: pip install python-pptx")
        
        prs = Presentation(str(path))
        parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = ''.join(run.text for run in para.runs)
                        if text.strip():
                            parts.append(text)
            parts.append('')
        
        return '\n'.join(parts)

    @staticmethod
    def extract_segments(path):
        """Slide-based segmentation."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("PowerPoint reading requires python-pptx: pip install python-pptx")
        
        prs = Presentation(str(path))
        segments = []
        
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = ''.join(run.text for run in para.runs)
                        if text.strip():
                            texts.append(text)
            
            segment_text = '\n'.join(texts)
            if segment_text:
                segments.append({
                    'text': segment_text,
                    'start': i - 1,
                    'end': i,
                    'tokens': len(segment_text) // 4,
                    'label': f'Slide {i}'
                })
        
        return segments


reader_registry.register(['.pptx'], PptxReader, requires=['pptx'])
